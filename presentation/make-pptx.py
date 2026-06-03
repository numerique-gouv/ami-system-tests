#!/usr/bin/env python3
"""
Assemble les PNG d'un dossier en un fichier PPTX (1 image par slide, 16:9).

Usage :
  python3 make-pptx.py <dossier-png> <output.pptx>
  python3 make-pptx.py --video-on-slide 4 video.mp4 <dossier-png> <output.pptx>

L'option --video-on-slide remplace le PNG de la slide N (base 1) par la vidéo MP4
avec son thumbnail extrait via ffmpeg (si disponible).
"""
import argparse
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


def extract_thumbnail(video_path: Path) -> Path | None:
    fd, tmp_str = tempfile.mkstemp(suffix=".png")
    import os; os.close(fd)
    tmp = Path(tmp_str)
    result = subprocess.run(
        ["ffmpeg", "-i", str(video_path), "-vframes", "1", "-q:v", "2", str(tmp), "-y"],
        capture_output=True,
    )
    return tmp if result.returncode == 0 else None


def main():
    parser = argparse.ArgumentParser(
        description="Assemble des PNG en PPTX 16:9, avec vidéo optionnelle sur une slide."
    )
    parser.add_argument(
        "--video-on-slide",
        nargs=2,
        metavar=("N", "VIDEO"),
        help="Numéro de slide (1-based) et chemin vers la vidéo MP4 à embarquer.",
    )
    parser.add_argument("png_dir", help="Dossier contenant les PNG (triés par nom)")
    parser.add_argument("output", help="Fichier PPTX de sortie")
    args = parser.parse_args()

    png_dir = Path(args.png_dir)
    output = Path(args.output)

    pngs = sorted(png_dir.glob("*.png"))
    if not pngs:
        print(f"Aucun PNG trouvé dans {png_dir}")
        raise SystemExit(1)

    video_slide_num: int | None = None
    video_path: Path | None = None
    if args.video_on_slide:
        video_slide_num = int(args.video_on_slide[0])
        video_path = Path(args.video_on_slide[1])
        if not video_path.exists():
            print(f"Vidéo introuvable : {video_path}")
            raise SystemExit(1)

    # 16:9 en EMU (914400 EMU = 1 inch, 10 × 5.625 inches)
    width = Emu(9144000)
    height = Emu(5143500)

    prs = Presentation()
    prs.slide_width = width
    prs.slide_height = height
    blank_layout = prs.slide_layouts[6]

    thumbnail: Path | None = None

    for i, png in enumerate(pngs, start=1):
        slide = prs.slides.add_slide(blank_layout)

        if video_slide_num and i == video_slide_num and video_path:
            # Slide vidéo : thumbnail plein écran en fond + vidéo embarquée par-dessus
            if thumbnail is None:
                thumbnail = extract_thumbnail(video_path)
                if thumbnail:
                    print(f"🎞  Thumbnail extrait : {thumbnail}")
                else:
                    print("⚠️  ffmpeg indisponible — slide sans thumbnail.")

            bg_image = thumbnail if thumbnail else png
            slide.shapes.add_picture(str(bg_image), Emu(0), Emu(0), width, height)
            slide.shapes.add_movie(
                str(video_path),
                Emu(0),
                Emu(0),
                width,
                height,
                poster_frame_image=str(thumbnail) if thumbnail else None,
                mime_type="video/mp4",
            )
            print(f"🎬 Vidéo embarquée sur la slide {i}")
        else:
            slide.shapes.add_picture(str(png), Emu(0), Emu(0), width, height)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(f"✅ {len(pngs)} slides → {output}")


if __name__ == "__main__":
    main()
